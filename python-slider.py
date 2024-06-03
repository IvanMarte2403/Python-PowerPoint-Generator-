from pptx import Presentation
from pptx.util import Pt
import openai
from config import api_key  
from grafics import add_header_bar, download_image
from pptx.util import Inches, Pt
from promts import introduccionPrompt, tituloClase, prompts

openai.api_key = api_key    


#Función para generar texto con GPT-3
def generate_text(prompt):

    print(f"Prompt: {prompt}")

    messages = [
       
        {
            "role" : "user", "content" : prompt     
        },
    ]


    response = openai.chat.completions.create(
        model="gpt-3.5-turbo",  # Selección  de Módelo de Chat
        messages= messages,
        temperature= 0
    )
    return response.choices[0].message.content



#Función para generar imagen con DALL-E
def generate_dalle_image(prompt):
    response = openai.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size="1024x1024",
        quality="standard",
        n=1,  # Generate a single image
    )
    image_url = response.data[0].url
    return image_url


#=======================================================================================================
#===============================[ Creación de la presentación en PowerPoint ]===========================


#Función para crear presentación en PowerPoint

def create_presentation(prompts, file_name="presentation.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.2)

    #================================[Diapositivas de la presentación]===================================
    # Diapositiva de inicio
    slide = prs.slides.add_slide(prs.slide_layouts[1])

     # Añadir imagen al fondo de la diapositiva
    slide.shapes.add_picture('img/Background.png', Inches(0), prs.slide_height - Inches(5), prs.slide_width, Inches(5))
    #Añadir imagen de 'DataRebels' nombre 
    slide.shapes.add_picture('img/logo-datarebels-rosa.png', Inches(0.5), Inches(1.5), Inches(3), Inches(0.7))
    #Añadir iamgen de 'DataRebels' logo
    slide.shapes.add_picture('img/logo-datarebels-portada.png', Inches(10), Inches(2), Inches(4), Inches(4))    
    # Añadir título de la presentación
    title_shape = slide.shapes.title
    title_shape.text = tituloClase

    # Cambiar la tipografía y el tamaño del texto
    title_shape.text_frame.paragraphs[0].runs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(45) 
    title_shape.left = Inches(0.5)
    title_shape.top = Inches(3)
    # Aumentar el ancho del cuadro de texto del título
    title_shape.width = Inches(7)
    # Mover el título al frente
    title_shape._element.getparent().append(title_shape._element)

   

    #======================================================================================================

    #Agrega una diapositiva 
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Objetivos"
    #Edición de los objetivos de la clase
    slide.placeholders[1].text = generate_text("Objetivos de la clase")

    # Diapositivas de temas
    for i, prompt in enumerate(prompts, start=1):
        # Diapositiva del tema
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Tema {i}"
        slide.placeholders[1].text = generate_text(prompt)
    
       
        #Add Grafic Bar 
        add_header_bar(slide, "FF0E68")

        # Diapositiva del ejemplo o analogía
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Tema {i} - Ejemplo o analogía"
        #Genera el ejemplo o la analogía 
        slide.placeholders[1].text = generate_text(f"Ejemplo o analogía para {prompt}")

        add_header_bar(slide, "FF0E68")


        # Generate DALL-E image for analogy
        analogy_prompt = f"Ejemplo o analogía visual para {prompt}"
        analogy_image_url = generate_dalle_image(analogy_prompt)

        # Download DALL-E image
        image_filename = f"analogy_image_{i}.png"
        download_image(analogy_image_url, image_filename)

         # Add DALL-E image to analogy slide (after text)

        slide.shapes.add_picture(
            image_filename,
            left=prs.slide_width / 4,
            top=prs.slide_height / 3,
            width=prs.slide_width / 2,
            height=prs.slide_height / 2,
        )  # Adjust positioning as needed
        
# Ajustar el tamaño del texto
    for slide in prs.slides:
       for shape in slide.shapes:
        # Verificar si la forma es el título de la diapositiva
        if shape == slide.shapes.title:
            continue  # Si es el título, saltar al siguiente ciclo

        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(16)
                    run.font.name = 'Open Sans'


    prs.save(file_name)
    print(f"Presentation saved as {file_name}")


#Función principal

def main():
    

    # Generar texto para cada prompt
    generated_texts = [generate_text(prompt) for prompt in prompts]

    # Crear la presentación
    create_presentation(generated_texts)

if __name__ == "__main__":
    main()



    